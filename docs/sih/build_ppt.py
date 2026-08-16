#!/usr/bin/env python3
"""Build the SIH 2026 6-slide idea PPT for STS Safety from the official template."""

from __future__ import annotations

import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
TEMPLATE = Path(
    "/Users/mitansh7/Downloads/SIH2026-IDEA-Presentation-Format (1).pptx"
)
OUT_PPTX = ROOT / "STS-Safety-SIH2026-Idea-Presentation.pptx"
ASSETS = ROOT / "assets"

# SIH palette (saffron / navy / green from official logo)
NAVY = RGBColor(0x0B, 0x25, 0x45)
NAVY2 = RGBColor(0x12, 0x33, 0x5A)
ORANGE = RGBColor(0xF3, 0x6C, 0x21)
GREEN = RGBColor(0x1B, 0x8A, 0x4A)
BLUE = RGBColor(0x00, 0x70, 0xC0)
TEAL = RGBColor(0x0E, 0x6B, 0x6B)
RED = RGBColor(0xB4, 0x23, 0x18)
GOLD = RGBColor(0xC4, 0x8A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF5, 0xF7, 0xFA)
LIGHT = RGBColor(0xE8, 0xEE, 0xF5)
MUTED = RGBColor(0x4A, 0x55, 0x63)
INK = RGBColor(0x1A, 0x22, 0x2E)
SOFT = RGBColor(0xD6, 0xDE, 0xE8)

TEAM = "STS Safety"
PRODUCT = "STS SAFETY"

EMU_IN = 914400


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor | None, pt: float = 0.75) -> None:
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def font(run, name="Calibri", size=10, bold=False, color=INK, italic=False) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    # Force latin + ea so PowerPoint doesn't swap to Calibri Light on Mac
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", name)


def _tf(shape, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    return tf


def add_shape(slide, st, l, t, w, h, fill=None, line=None, line_pt=0.75):
    sh = slide.shapes.add_shape(st, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is not None:
        set_fill(sh, fill)
    else:
        sh.fill.background()
    set_line(sh, line, line_pt)
    # kill default text margin bloat
    tf = sh.text_frame
    tf.word_wrap = True
    for attr, val in (("lIns", 72000), ("rIns", 72000), ("tIns", 36000), ("bIns", 36000)):
        tf._txBody.bodyPr.set(attr, str(val))
    return sh


def rect(slide, l, t, w, h, fill, line=None, line_pt=0.75):
    return add_shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, fill, line, line_pt)


def round_rect(slide, l, t, w, h, fill, line=None, line_pt=0.75, adj=0.08):
    sh = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill, line, line_pt)
    try:
        sh.adjustments[0] = adj
    except Exception:
        pass
    return sh


def oval(slide, l, t, w, h, fill, line=None):
    return add_shape(slide, MSO_SHAPE.OVAL, l, t, w, h, fill, line)


def tb(slide, l, t, w, h, text, size=10, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, name="Calibri"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = _tf(box, anchor)
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    r = p.add_run()
    r.text = text
    font(r, name=name, size=size, bold=bold, color=color, italic=italic)
    return box


def lines(slide, l, t, w, h, items, anchor=MSO_ANCHOR.TOP, valign_center=False):
    """items: list of dict(text, size, bold, color, align, italic, space_after)."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = _tf(box, MSO_ANCHOR.MIDDLE if valign_center else anchor)
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = it.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(it.get("sb", 0))
        p.space_after = Pt(it.get("sa", 2))
        p.level = it.get("level", 0)
        r = p.add_run()
        r.text = it["text"]
        font(
            r,
            name=it.get("name", "Calibri"),
            size=it.get("size", 10),
            bold=it.get("bold", False),
            color=it.get("color", INK),
            italic=it.get("italic", False),
        )
    return box


def card_title_body(slide, l, t, w, h, title, body_items, fill=WHITE, accent=ORANGE, title_color=None):
    sh = round_rect(slide, l, t, w, h, fill, SOFT, 0.6, adj=0.06)
    rect(slide, l, t, 0.07, h, accent)
    tb(slide, l + 0.14, t + 0.06, w - 0.22, 0.28, title, 11, True, title_color or accent)
    lines(slide, l + 0.14, t + 0.34, w - 0.22, h - 0.42, body_items)
    return sh


def metric(slide, l, t, w, h, value, label, sub=None, fill=NAVY, value_color=ORANGE, label_color=WHITE):
    round_rect(slide, l, t, w, h, fill, None, adj=0.1)
    items = [
        {"text": value, "size": 16, "bold": True, "color": value_color, "align": PP_ALIGN.CENTER, "sa": 0},
        {"text": label, "size": 8, "bold": True, "color": label_color, "align": PP_ALIGN.CENTER, "sa": 0},
    ]
    if sub:
        items.append({"text": sub, "size": 7, "bold": False, "color": RGBColor(0xC5, 0xD0, 0xDC), "align": PP_ALIGN.CENTER, "sa": 0, "italic": True})
    lines(slide, l + 0.04, t + 0.08, w - 0.08, h - 0.12, items, valign_center=True)


def pill(slide, l, t, w, h, text, fill=ORANGE, color=WHITE, size=9):
    round_rect(slide, l, t, w, h, fill, None, adj=0.5)
    tb(slide, l, t, w, h, text, size, True, color, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def screenshot(slide, path: Path, l, t, w, h, caption: str) -> None:
    """Place a real product capture in a clean, labelled frame."""
    rect(slide, l - 0.04, t - 0.04, w + 0.08, h + 0.08, WHITE, SOFT, 0.8)
    slide.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))
    tb(slide, l, t + h + 0.06, w, 0.22, caption, 7.5, True, MUTED, PP_ALIGN.CENTER)


def delete_shape(shape) -> None:
    el = shape._element
    parent = el.getparent()
    if parent is not None:
        parent.remove(el)


def find_by_name(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_shape_text(shape, text: str, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    font(r, size=size, bold=bold, color=color)


def replace_textbox_lines(shape, items) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = it.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(it.get("sb", 0))
        p.space_after = Pt(it.get("sa", 4))
        r = p.add_run()
        r.text = it["text"]
        font(r, size=it.get("size", 14), bold=it.get("bold", False), color=it.get("color", INK))


def strip_instruction_box(slide) -> None:
    for sh in list(slide.shapes):
        if not sh.has_text_frame:
            continue
        blob = " ".join(p.text for p in sh.text_frame.paragraphs)
        needles = (
            "Proposed Solution (Describe",
            "Technologies to be used",
            "Analysis of the feasibility",
            "Potential impact on the target",
            "Details / Links of the reference",
            "Detailed explanation of the proposed",
        )
        if any(n in blob for n in needles):
            delete_shape(sh)


def set_footer_team(slide) -> None:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        t = (sh.text_frame.text or "").strip()
        if t in {"Your Team Name", "YourTeam Name", "Your Team Name "}:
            set_shape_text(sh, TEAM, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            # colour the oval
            try:
                set_fill(sh, ORANGE)
                set_line(sh, ORANGE)
            except Exception:
                pass


def set_title(slide, text: str, color=ORANGE, size=26) -> None:
    # Prefer placeholder title. Shift right so it does not collide with the
    # official "Team Name" oval at top-left of the SIH template.
    for sh in slide.shapes:
        if sh.has_text_frame and sh.name.startswith("Title"):
            sh.left = Inches(1.72)
            sh.top = Inches(0.08)
            sh.width = Inches(8.6)
            sh.height = Inches(0.72)
            tf = sh.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = text
            font(r, size=size, bold=True, color=color)
            return


def build() -> Path:
    if not TEMPLATE.exists():
        raise SystemExit(f"Missing template: {TEMPLATE}")
    shutil.copy(TEMPLATE, OUT_PPTX)
    prs = Presentation(str(OUT_PPTX))

    # ----- SLIDE 1: TITLE -----
    s1 = prs.slides[0]
    for sh in s1.shapes:
        if not sh.has_text_frame:
            continue
        joined = " ".join(p.text for p in sh.text_frame.paragraphs)
        if "SMART INDIA HACKATHON" in joined:
            replace_textbox_lines(
                sh,
                [
                    {"text": "SMART INDIA HACKATHON 2026", "size": 28, "bold": True, "color": NAVY, "sa": 2},
                ],
            )
        elif "TITLE PAGE" in joined:
            delete_shape(sh)
            continue
        elif "Problem" in joined and "Theme" in joined:
            sh.left = Inches(0.62)
            sh.top = Inches(1.48)
            sh.width = Inches(6.18)
            sh.height = Inches(4.95)
            replace_textbox_lines(
                sh,
                [
                    {"text": "Problem Statement ID  –  SIH25002", "size": 16, "bold": True, "color": NAVY, "sa": 8, "sb": 6},
                    {
                        "text": "Problem Statement Title  –  Smart Tourist Safety Monitoring & Incident Response System using AI, Geo-Fencing, and Blockchain-based Digital ID",
                        "size": 15,
                        "bold": False,
                        "color": INK,
                        "sa": 8,
                    },
                    {"text": "Theme  –  Travel & Tourism", "size": 16, "bold": False, "color": INK, "sa": 8},
                    {"text": "PS Category  –  Software", "size": 16, "bold": False, "color": INK, "sa": 8},
                    {
                        "text": "Organisation  –  Ministry of Development of North Eastern Region (MDoNER)",
                        "size": 15,
                        "bold": False,
                        "color": INK,
                        "sa": 8,
                    },
                    {"text": f"Team Name  –  {TEAM}", "size": 16, "bold": True, "color": ORANGE, "sa": 0},
                ],
            )
    screenshot(
        s1,
        ASSETS / "landing-crop.png",
        7.14,
        1.55,
        5.66,
        3.42,
        "Working prototype · tourist and command surfaces",
    )
    pill(s1, 8.60, 5.25, 2.74, 0.34, "LIVE PRODUCT · ₹0 INFRA", GREEN, WHITE, 10)

    # ----- SLIDES 2–6 content -----
    for idx in range(1, 6):
        strip_instruction_box(prs.slides[idx])
        set_footer_team(prs.slides[idx])

    fill_proposed(prs.slides[1])
    fill_technical(prs.slides[2])
    fill_feasibility(prs.slides[3])
    fill_impact(prs.slides[4])
    fill_research(prs.slides[5])

    # Drop the instructions slide (slide 7)
    rId = prs.slides._sldIdLst[-1].get(qn("r:id"))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])

    prs.save(str(OUT_PPTX))
    return OUT_PPTX


def fill_proposed(slide) -> None:
    set_title(slide, PRODUCT, ORANGE, 28)
    # official pointer chip
    pill(slide, 0.22, 0.92, 2.55, 0.26, "Proposed Solution", ORANGE, WHITE, 10)
    tb(
        slide,
        2.85,
        0.90,
        7.4,
        0.28,
        "Safety that works when the hills go dark  ·  Working prototype, not a mock-up",
        11,
        False,
        MUTED,
        italic=True,
    )

    # PROBLEM
    round_rect(slide, 0.18, 1.26, 4.22, 3.52, WHITE, SOFT, 0.7, 0.04)
    rect(slide, 0.18, 1.26, 4.22, 0.34, NAVY)
    tb(slide, 0.30, 1.28, 4.0, 0.30, "1.  The problem  —  North-East tourist safety", 11, True, WHITE)

    lines(
        slide,
        0.30,
        1.66,
        3.96,
        3.02,
        [
            {"text": "Tourism is national-scale. Coverage is not.", "size": 10, "bold": True, "color": ORANGE, "sa": 6},
            {
                "text": "•  1.21 crore domestic + 2.21 lakh foreign visits to the 8 NE states in 2023 (PIB / State Tourism Depts).",
                "size": 9,
                "color": INK,
                "sa": 5,
            },
            {
                "text": "•  NE share of India’s foreign tourist arrivals is only 2.05% (1.95 lakh; MoT Compendium 2024) — a priority region that is still under-visited.",
                "size": 9,
                "color": INK,
                "sa": 5,
            },
            {
                "text": "•  7 of 8 NE states have no dedicated Tourist Police. MoT’s 15-state list includes Sikkim only (Lok Sabha 27.07.2026).",
                "size": 9,
                "color": INK,
                "sa": 5,
            },
            {
                "text": "•  NCRB does not maintain crime data on domestic tourists. 156 cases against foreign tourists were registered in 2024 (NCRB / LS AU1299) — missing-person, restricted-zone and signal-loss events are invisible.",
                "size": 9,
                "color": INK,
                "sa": 5,
            },
            {
                "text": "•  Check-post, hotel, forest office and police do not share a database. A Tawang checkpoint cannot verify a Guwahati-issued ID if the issuing server is down.",
                "size": 9,
                "color": INK,
                "sa": 4,
            },
            {
                "text": "PS ask: digital ID + geo-fence + AI + SOS + E-FIR — with DPDP-safe privacy.",
                "size": 9,
                "bold": True,
                "color": NAVY,
                "sa": 0,
            },
        ],
    )

    # 3-LAYER SOLUTION
    round_rect(slide, 4.52, 1.26, 4.38, 3.52, WHITE, SOFT, 0.7, 0.04)
    rect(slide, 4.52, 1.26, 4.38, 0.34, ORANGE)
    tb(slide, 4.64, 1.28, 4.14, 0.30, "2.  Three-layer system  (the prototype)", 11, True, WHITE)

    layers = [
        ("01", "Geo-fence in Postgres", "AFTER INSERT trigger + PostGIS ST_Covers / ST_DWithin. The phone cannot skip an incident. Debounce via partial unique index so a border-stander does not generate 400 alerts."),
        ("02", "Soulbound digital ID", "ERC-5192 locked TDID on Polygon Amoy. Chain stores only keccak256(KYC ‖ salt) + itinerary hash. KYC stays pgcrypto-encrypted in Postgres. Selective disclosure at /verify."),
        ("03", "AI beside the path — never on it", "Tier-1 SQL rules always fire. IsolationForest only ranks. Groq/Gemini write the brief & E-FIR draft. If HF Space, Groq and Amoy are all down, SOS still lands."),
    ]
    y = 1.68
    for num, title, body in layers:
        oval(slide, 4.66, y, 0.32, 0.32, NAVY)
        tb(slide, 4.66, y + 0.03, 0.32, 0.26, num, 8, True, WHITE, PP_ALIGN.CENTER)
        tb(slide, 5.06, y - 0.02, 3.70, 0.22, title, 11, True, NAVY)
        tb(slide, 5.06, y + 0.22, 3.70, 0.72, body, 8.5, False, MUTED)
        y += 0.98

    # UNIQUENESS
    round_rect(slide, 9.02, 1.26, 4.12, 3.52, NAVY, None, adj=0.04)
    tb(slide, 9.16, 1.34, 3.84, 0.28, "3.  What no typical build does", 12, True, ORANGE)
    uniq = [
        ("Safety path is the database.", "Blockchain and AI sit beside it. Unplug the WAN: Realtime + Anvil + ONNX + PMTiles still demo."),
        ("On-device warning first.", "Turf.js evaluates cached zones before the packet arrives. IndexedDB + Background Sync + sms: fallback for Arunachal towers."),
        ("Lawful identity, not an NFT.", "Zero PII on-chain. Hotel learns match / no-match, never the passport number. Expired trip auto-invalidates on-chain."),
        ("Evidentiary E-FIR, not a chatbot FIR.", "Officer must approve. PDF hash is anchored. LLM never files to CCTNS."),
        ("Honest ML.", "Trained on synthetic NE trajectories (DPDP). Hold-out: 0 FP on 1,500 normal-trek windows. Rules catch what the forest misses."),
    ]
    y = 1.68
    for t, b in uniq:
        tb(slide, 9.16, y, 3.84, 0.20, "▸  " + t, 10, True, WHITE)
        tb(slide, 9.32, y + 0.20, 3.68, 0.42, b, 8, False, RGBColor(0xC9, 0xD4, 0xE0))
        y += 0.60

    # BOTTOM METRICS
    metric(slide, 0.18, 4.88, 2.52, 0.92, "180–450 ms", "SOS → first channel", "measured, local stack", NAVY, ORANGE)
    metric(slide, 2.82, 4.88, 2.52, 0.92, "0 FP / 1,500", "normal-trek hold-out", "IsolationForest; rules still decide", TEAL, WHITE, WHITE)
    metric(slide, 5.46, 4.88, 2.52, 0.92, "₹0 infra", "no credit card in the stack", "Vercel + Supabase + HF + Amoy", GREEN, WHITE, WHITE)
    metric(slide, 8.10, 4.88, 2.52, 0.92, "5 locales", "en · hi · as · bn · ne", "PS multilingual requirement", NAVY2, ORANGE)
    metric(slide, 10.74, 4.88, 2.40, 0.92, "< 2 s", "Telegram target (online)", "panic path skips debounce", ORANGE, WHITE, WHITE)

    tb(
        slide,
        0.20,
        5.84,
        12.9,
        0.22,
        "Sources (problem stats): PIB ‘Tourism in North-Eastern Region’ Dec 2024; MoT India Tourism Data Compendium Key Highlights 2024; Lok Sabha USQ 1299, 27.07.2026 (NCRB). Product stats: repo README + services/ai/artifacts/metrics.json.",
        8,
        False,
        MUTED,
        italic=True,
    )


def fill_technical(slide) -> None:
    set_title(slide, "TECHNICAL APPROACH", ORANGE, 24)
    pill(slide, 8.55, 0.22, 1.85, 0.26, "Working prototype", GREEN, WHITE, 9)

    # stack table header
    tb(slide, 0.20, 0.95, 6.4, 0.24, "Technologies  —  every row is free-tier, no card", 12, True, NAVY)

    headers = ["Layer", "Choice", "Why this, not the obvious"]
    rows = [
        ["App", "Next.js 16 PWA + MapLibre", "URL, not an APK. Web Push + Background Sync included."],
        ["Hot path", "PostgREST insert + RLS", "Skips serverless. 10 writes/s is noise for Postgres."],
        ["Geo engine", "PostGIS 3  ·  plpgsql trigger", "ST_Covers in the database. Client cannot bypass."],
        ["Maps", "OpenFreeMap + Turf.js", "No Mapbox/Google key. On-device pre-check."],
        ["Identity", "Solidity 0.8.24  ·  ERC-5192", "Soulbound TDID. keccak256 commitment only."],
        ["Chain", "Polygon Amoy / Anvil", "Free POL. Offline = Foundry Anvil."],
        ["ML", "IsolationForest → ONNX", "HF Space + onnxruntime-node fallback."],
        ["LLM", "Groq → Gemini → template", "Brief + E-FIR draft. Never gates the alert."],
        ["Notify", "Realtime + VAPID + Telegram", "Twilio/SMS is a stub for a ministry short-code."],
        ["Privacy", "pgcrypto + RLS + 24h purge", "DPDP: consent, purpose-limit, right to erasure."],
    ]

    x0, y0, widths = 0.18, 1.22, [1.15, 2.35, 3.05]
    # header
    x = x0
    for i, h in enumerate(headers):
        rect(slide, x, y0, widths[i], 0.28, NAVY)
        tb(slide, x + 0.04, y0 + 0.02, widths[i] - 0.06, 0.24, h, 9, True, WHITE)
        x += widths[i]
    for r_i, row in enumerate(rows):
        y = y0 + 0.28 + r_i * 0.33
        bg = OFFWHITE if r_i % 2 == 0 else WHITE
        x = x0
        for i, cell in enumerate(row):
            rect(slide, x, y, widths[i], 0.33, bg, SOFT, 0.4)
            tb(
                slide,
                x + 0.05,
                y + 0.04,
                widths[i] - 0.08,
                0.26,
                cell,
                8 if i == 2 else 9,
                i == 0,
                ORANGE if i == 0 else INK,
            )
            x += widths[i]

    # methodology flow
    tb(slide, 6.85, 0.95, 6.3, 0.24, "Methodology  —  SOS / geo-fence hot path", 12, True, NAVY)

    flow = [
        ("1", "Device", "watchPosition\nTurf.js local zone\nIndexedDB if offline"),
        ("2", "Ingest", "RLS insert into\nlocation_pings\n~80–150 ms"),
        ("3", "PostGIS", "evaluate_position()\nzone · corridor ·\nspeed · silence"),
        ("4", "Incident", "INSERT incidents\nRealtime <100 ms\ncontrol-room siren"),
        ("5", "Side path", "ONNX score\nLLM brief\nchain anchor"),
        ("6", "Dispatch", "nearest unit\nTelegram + Push\nE-FIR draft"),
    ]
    fx, fy = 6.85, 1.26
    fw = 0.95
    for i, (n, t, b) in enumerate(flow):
        x = fx + (i % 3) * 2.10
        y = fy + (i // 3) * 1.55
        round_rect(slide, x, y, 1.92, 1.38, WHITE, SOFT, 0.7, 0.08)
        oval(slide, x + 0.08, y + 0.10, 0.28, 0.28, ORANGE)
        tb(slide, x + 0.08, y + 0.12, 0.28, 0.24, n, 10, True, WHITE, PP_ALIGN.CENTER)
        tb(slide, x + 0.40, y + 0.12, 1.42, 0.24, t, 12, True, NAVY)
        tb(slide, x + 0.10, y + 0.44, 1.72, 0.86, b, 9, False, MUTED)

    # SOS callout
    round_rect(slide, 6.85, 4.42, 3.42, 1.42, NAVY, None, adj=0.05)
    tb(slide, 7.00, 4.50, 3.14, 0.24, "Panic path  —  button to siren", 11, True, ORANGE)
    tb(
        slide,
        7.00,
        4.76,
        3.14,
        0.98,
        "Hold 1.5 s → direct RLS incident insert → Realtime takeover → Telegram + Push + email in parallel.\nOffline: Background Sync + pre-composed SMS with coordinates.\nAI and blockchain stay beside the safety path.",
        8.5,
        False,
        WHITE,
    )
    screenshot(
        slide,
        ASSETS / "tourist-login-crop.png",
        10.45,
        4.42,
        1.22,
        1.07,
        "Tourist PWA",
    )
    screenshot(
        slide,
        ASSETS / "officer-login-crop.png",
        11.85,
        4.42,
        1.22,
        1.03,
        "Command centre",
    )

    tb(
        slide,
        0.20,
        6.18,
        6.4,
        0.22,
        "Surfaces: /home  /sos  /onboard  /dashboard  /verify  ·  Simulator stands in for the IoT band on the same ingest path.",
        8,
        False,
        MUTED,
        italic=True,
    )


def fill_feasibility(slide) -> None:
    set_title(slide, "FEASIBILITY AND VIABILITY", ORANGE, 22)

    # POC
    round_rect(slide, 0.18, 1.02, 6.55, 2.55, WHITE, SOFT, 0.7, 0.04)
    rect(slide, 0.18, 1.02, 6.55, 0.32, GREEN)
    tb(slide, 0.32, 1.04, 6.30, 0.28, "Feasibility  —  already built and measured", 12, True, WHITE)

    poc = [
        ("Hold-out IsolationForest", "P = 1.00   R = 0.56   F1 = 0.72", "n_test = 1,650  ·  threshold 0.72  ·  18 features"),
        ("Confusion (hold-out)", "TN 1,500  ·  FP 0  ·  TP 84  ·  FN 66", "0 false positives on 1,500 normal-trek windows"),
        ("Scenario recall", "zone-breach 1.00  ·  stationary-anomaly 1.00", "route-deviation 0.00 — caught by PostGIS rules, not the forest"),
        ("SOS latency", "180–450 ms first channel (local)", "Telegram target < 2,000 ms when the bot is live"),
        ("Cost / scale", "₹0 at demo  ·  24 h ping downsample", "50 tourists × 1 ping / 5 s ≈ 104 MB/day before downsample"),
    ]
    y = 1.40
    for a, b, c in poc:
        tb(slide, 0.32, y, 1.85, 0.38, a, 9, True, NAVY)
        tb(slide, 2.20, y, 2.55, 0.22, b, 9, True, ORANGE)
        tb(slide, 2.20, y + 0.18, 4.30, 0.20, c, 8, False, MUTED)
        y += 0.42

    # viability
    round_rect(slide, 6.86, 1.02, 6.28, 2.55, WHITE, SOFT, 0.7, 0.04)
    rect(slide, 6.86, 1.02, 6.28, 0.32, ORANGE)
    tb(slide, 7.00, 1.04, 6.00, 0.28, "Viability  —  government adoption, not a startup deck", 12, True, WHITE)
    lines(
        slide,
        7.00,
        1.42,
        6.00,
        2.05,
        [
            {"text": "Who pays: MDoNER / state tourism + police. We are a processor; they are the data fiduciary.", "size": 9, "color": INK, "sa": 5},
            {"text": "Why they would: 7 NE states have no Tourist Police. This is a force-multiplier for 112 / control rooms, not a consumer app.", "size": 9, "color": INK, "sa": 5},
            {"text": "Month-one production (if funded): (1) ministry SMS short-code, (2) real ILP/KYC adapter, (3) non-pausing Postgres, (4) Amoy → Polygon mainnet with a hardware-backed issuer, (5) CERT-In threat model. The geofence engine does not change.", "size": 9, "color": INK, "sa": 5},
            {"text": "We do not claim a ₹-crore TAM or a 20% crime drop. Those numbers are not in any official series we can cite.", "size": 9, "bold": True, "color": NAVY, "sa": 0},
        ],
    )

    # challenges
    tb(slide, 0.20, 3.66, 13.0, 0.24, "Challenges, risks, and the strategy already in the repo", 12, True, NAVY)

    challenges = [
        ("01", "Low connectivity\n(Arunachal / hills)", "IndexedDB queue + Background Sync. Turf.js warns on-device. SOS falls back to sms: with coordinates. Offline demo: four env switches (DB / chain / AI / tiles)."),
        ("02", "GPS spoof / rooted\nphone", "Cannot stop a rooted device. Refuse implausible speed (>150 km/h). Session JWT. SOS is higher-trust because it is a user gesture. Checkpoint QR still binds the person."),
        ("03", "Synthetic ML data\n(DPDP)", "Stated in the pitch. IsolationForest is unsupervised on generated NE tracks. The decision to alert is rules + PostGIS. Rules-only is one env flip."),
        ("04", "LLM + legal E-FIR", "Correct: an LLM cannot file. We draft a narrative PDF. approved_by is a human. Nothing is submitted to CCTNS from this system."),
        ("05", "Free-tier pause /\nrate-limit on stage", "GitHub Action hits /api/health every 6 h. Local supabase start + Anvil + ONNX + PMTiles is the stage fallback. Rehearsed unplugged."),
        ("06", "Privacy / DPDP /\n“is this surveillance?”", "Opt-in, trip-scoped, purpose = safety. KYC encrypted, never on-chain, never logged. RLS on every table. 24 h ping retention. Emergency contacts only at critical."),
    ]
    for i, (n, t, b) in enumerate(challenges):
        col, row = i % 3, i // 3
        x = 0.18 + col * 4.38
        y = 3.94 + row * 1.18
        round_rect(slide, x, y, 4.22, 1.10, WHITE, SOFT, 0.6, 0.07)
        oval(slide, x + 0.10, y + 0.12, 0.36, 0.36, NAVY)
        tb(slide, x + 0.10, y + 0.16, 0.36, 0.28, n, 9, True, WHITE, PP_ALIGN.CENTER)
        tb(slide, x + 0.52, y + 0.10, 3.58, 0.40, t, 11, True, ORANGE)
        tb(slide, x + 0.12, y + 0.50, 3.98, 0.54, b, 8, False, MUTED)


def fill_impact(slide) -> None:
    set_title(slide, "IMPACT AND BENEFITS", ORANGE, 24)

    # official stats strip
    stats = [
        ("9.95 million", "Foreign tourist arrivals, 2024", "MoT / BoI  ·  PIB"),
        ("2.51 billion", "Domestic tourist visits, 2023", "MoT  ·  PIB"),
        ("5.22%", "Tourism share of GDP, 2023-24", "TSA / PLFS, provisional"),
        ("8.46 crore", "Tourism jobs (direct+indirect)", "PLFS 2023-24"),
        ("1.21 crore", "NE domestic visits, 2023", "PIB NER  ·  State Depts"),
    ]
    for i, (v, l, s) in enumerate(stats):
        x = 0.18 + i * 2.62
        round_rect(slide, x, 1.00, 2.52, 1.12, NAVY if i % 2 == 0 else NAVY2, None, adj=0.08)
        tb(slide, x + 0.08, 1.08, 2.36, 0.36, v, 16, True, ORANGE, PP_ALIGN.CENTER)
        tb(slide, x + 0.08, 1.44, 2.36, 0.36, l, 9, True, WHITE, PP_ALIGN.CENTER)
        tb(slide, x + 0.08, 1.78, 2.36, 0.24, s, 7, False, RGBColor(0xB8, 0xC4, 0xD2), PP_ALIGN.CENTER)

    # three benefit columns
    cols = [
        (
            "Social  —  who is safer",
            GREEN,
            [
                "Tourists (domestic + foreign) in Kaziranga, Tawang, Sohra, Loktak, border approaches — restricted-zone warning before they are lost.",
                "Women / solo travellers: hold-to-SOS that actually reaches a human, with last-known point and digital ID.",
                "Families: opt-in tracking, emergency contacts only at critical — not a live broadcast.",
                "Check-post / hotel / forest staff: one QR, green/amber/red, no central API dependency.",
                "5 UI languages (en, hi, as, bn, ne) matching the PS multilingual clause.",
            ],
        ),
        (
            "Economic  —  what it unlocks",
            ORANGE,
            [
                "Tourism is 5.22% of GDP and 84.63 million jobs. NE is a MoT + MDoNER priority that still takes only 2.05% of FTAs.",
                "Perception of safety is a binding constraint on high-value / foreign itineraries into the hills.",
                "Control-room MTTA/MTTR becomes measurable. Today most NE states have no Tourist Police unit to even record that clock.",
                "₹0 to run the prototype. First paid line in production is SMS + a non-pausing database — not a rewrite.",
                "E-FIR drafts cut officer paperwork; the officer still signs.",
            ],
        ),
        (
            "Governance / environment",
            BLUE,
            [
                "MDoNER PS delivered as specified: ID + geo-fence + AI + dashboard + E-FIR.",
                "Aligns with Tourist Police Scheme (BPR&D / MHA / MoT national conference) — a software layer those 7 NE states can actually use.",
                "DPDP 2023: data minimisation (no PII on-chain), purpose limitation, retention cron, right to erasure.",
                "Kaziranga / forest restricted polygons reduce illegal entry into core habitat — a conservation side-effect of the geo-fence, not a claimed ‘carbon’ product.",
                "Inner Line / LAC: we seed public caution polygons (Bum La approach), not classified overlays.",
            ],
        ),
    ]
    for i, (title, accent, bullets) in enumerate(cols):
        x = 0.18 + i * 4.38
        round_rect(slide, x, 2.24, 4.22, 3.55, WHITE, SOFT, 0.7, 0.04)
        rect(slide, x, 2.24, 4.22, 0.36, accent)
        tb(slide, x + 0.12, 2.28, 3.98, 0.28, title, 12, True, WHITE)
        y = 2.70
        for b in bullets:
            tb(slide, x + 0.14, y, 3.94, 0.58, "•  " + b, 9, False, INK)
            y += 0.60

    tb(
        slide,
        0.20,
        5.86,
        12.9,
        0.28,
        "We do not invent a ‘lives saved’ or ‘₹ X crore unlocked’ figure. Impact is: faster, auditable incident response for a region that currently has neither a tourist-police layer nor a domestic-tourist crime series.",
        9,
        True,
        NAVY,
    )


def fill_research(slide) -> None:
    set_title(slide, "RESEARCH AND REFERENCES", ORANGE, 22)

    # comparison table
    tb(slide, 0.20, 0.95, 8.4, 0.22, "Comparison with existing / typical SIH tourist-safety builds", 12, True, NAVY)

    headers = ["Capability", "STS Safety (this prototype)", "Typical conventional build"]
    rows = [
        ["Geo-fence", "PostGIS trigger; client cannot skip", "App if-else or Mongo $geoWithin"],
        ["Digital ID", "ERC-5192 soulbound + keccak256 commitment", "QR in a database, or NFT with PII"],
        ["AI role", "Ranks only. Rules decide. LLM never gates SOS", "LLM on GPS / ‘predictive hotspot’"],
        ["Offline", "Turf.js + IndexedDB + sms: + Anvil + ONNX", "Needs internet / paid Map SDK"],
        ["Privacy", "RLS, pgcrypto, 24 h ping purge, no PII on-chain", "Central Aadhaar log"],
        ["E-FIR", "Officer-approved PDF; hash anchored", "Chatbot ‘files’ a legal document"],
        ["Cost", "₹0, no card, four local fallbacks", "Twilio + Mapbox + Firebase Blaze"],
        ["ML honesty", "Synthetic tracks; metrics published", "Unstated accuracy on unseen GPS"],
    ]
    widths = [1.55, 5.55, 5.95]
    x0, y0 = 0.18, 1.20
    x = x0
    for i, h in enumerate(headers):
        rect(slide, x, y0, widths[i], 0.28, NAVY if i != 1 else ORANGE)
        tb(slide, x + 0.05, y0 + 0.02, widths[i] - 0.08, 0.24, h, 9, True, WHITE)
        x += widths[i]
    for r_i, row in enumerate(rows):
        y = y0 + 0.28 + r_i * 0.30
        bg = OFFWHITE if r_i % 2 == 0 else WHITE
        x = x0
        for i, cell in enumerate(row):
            rect(slide, x, y, widths[i], 0.30, bg, SOFT, 0.35)
            tb(slide, x + 0.05, y + 0.04, widths[i] - 0.08, 0.24, cell, 8, i == 0, NAVY if i == 0 else INK)
            x += widths[i]

    # future + refs
    round_rect(slide, 0.18, 4.00, 6.55, 1.88, WHITE, SOFT, 0.7, 0.04)
    rect(slide, 0.18, 4.00, 6.55, 0.30, TEAL)
    tb(slide, 0.32, 4.02, 6.30, 0.26, "Future scope  —  only what the architecture already allows", 11, True, WHITE)
    lines(
        slide,
        0.32,
        4.36,
        6.28,
        1.44,
        [
            {"text": "•  Ministry SMS short-code on the existing INotificationChannel adapter (sms.stub.ts).", "size": 9, "color": INK, "sa": 3},
            {"text": "•  Real Inner Line Permit / KYC adapter; custodial HD wallet → ERC-4337 paymaster.", "size": 9, "color": INK, "sa": 3},
            {"text": "•  Polygon mainnet + hardware-backed issuer. Contracts are vanilla OZ v5 — config, not a rewrite.", "size": 9, "color": INK, "sa": 3},
            {"text": "•  Physical IoT band posting to the same location_pings ingest (the simulator is that interface today).", "size": 9, "color": INK, "sa": 3},
            {"text": "•  CERT-In intake + threat model. Geofence engine stays in Postgres.", "size": 9, "color": INK, "sa": 0},
        ],
    )

    round_rect(slide, 6.86, 4.00, 6.28, 1.88, WHITE, SOFT, 0.7, 0.04)
    rect(slide, 6.86, 4.00, 6.28, 0.30, NAVY)
    tb(slide, 7.00, 4.02, 6.00, 0.26, "References  —  clickable in the PDF", 11, True, WHITE)
    refs = [
        (
            "MoT / PIB — Tourism in NER, DTV+FTV 2023",
            "https://tourism.gov.in/sites/default/files/2024-12/PIB2082598.pdf",
        ),
        (
            "MoT Compendium Key Highlights 2024 (NE FTA 2.05%)",
            "https://tourism.gov.in/sites/default/files/2025-02/India%20Tourism%20Data%20Compendium%20key%20highlights%202024.pdf",
        ),
        (
            "PIB — FTA 9.95M (2024); GDP 5.22%; jobs 84.63M",
            "https://www.pib.gov.in/PressReleasePage.aspx?PRID=2240657",
        ),
        (
            "Lok Sabha USQ 1299 (27.07.2026) — Tourist Police + NCRB 156",
            "https://sansad.in/getFile/lsapps/loksabhaquestions/annex/188/AU1299_510Xy0.pdf",
        ),
        ("DPDP Act 2023  ·  ERC-5192  ·  W3C Verifiable Credentials", None),
        ("Liu et al., Isolation Forest, ICDM 2008  ·  PostGIS ST_Covers / GiST", None),
        ("SIH PS SIH25002 — MDoNER, Travel & Tourism, Software", None),
    ]
    box = slide.shapes.add_textbox(Inches(7.00), Inches(4.34), Inches(6.00), Inches(1.48))
    tf = _tf(box)
    tf.clear()
    for i, (label, url) in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = "•  " + label
        font(r, size=8, color=BLUE if url else INK)
        if url:
            r.hyperlink.address = url

    tb(
        slide,
        0.20,
        5.92,
        12.9,
        0.22,
        "Prototype: tourist PWA + command centre + checkpoint verify + PostGIS engine + ONNX model + Anvil/Amoy contracts. Demo script: docs/DEMO-SCRIPT.md  ·  Hostile Q&A: docs/JUDGE-QA.md",
        8,
        True,
        MUTED,
        italic=True,
    )


if __name__ == "__main__":
    path = build()
    print(path)
